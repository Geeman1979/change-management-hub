"""Document generation utilities — Word documents and PowerPoint presentations."""

import os
import json
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt, Emu as PptEmu
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN


def _hex_to_rgb(hex_colour):
    """Convert hex colour string to RGB tuple."""
    h = hex_colour.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def _hex_to_pptrgb(hex_colour):
    """Convert hex colour to python-pptx RGBColor."""
    r, g, b = _hex_to_rgb(hex_colour)
    return PptRGBColor(r, g, b)


def generate_word_doc(communication_data, brand, output_path):
    """Generate a Word document from communication data with brand styling."""
    doc = Document()

    primary = _hex_to_rgb('#001D38')  # Gold Fields navy
    accent = _hex_to_rgb('#C8A064')   # Gold Fields gold
    secondary = _hex_to_rgb('#0a2744')
    teal = _hex_to_rgb('#00B398')

    # -- Default style tweaks --
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)
    font.color.rgb = RGBColor(*primary)

    # -- Header section --
    brand_name = brand.get('brand_name', 'Change Management Hub')
    header_para = doc.add_paragraph()
    header_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = header_para.add_run(brand_name.upper())
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(*accent)
    run.font.bold = True

    # Divider line
    divider = doc.add_paragraph()
    run = divider.add_run('_' * 70)
    run.font.color.rgb = RGBColor(*accent)
    run.font.size = Pt(6)

    # -- Title --
    subject = communication_data.get('subject', 'Communication')
    title = doc.add_heading(subject, level=1)
    for run in title.runs:
        run.font.color.rgb = RGBColor(*primary)
        run.font.size = Pt(18)

    # -- Greeting --
    greeting = communication_data.get('greeting', '')
    if greeting:
        p = doc.add_paragraph()
        run = p.add_run(greeting)
        run.font.size = Pt(11)
        run.font.bold = True

    # -- Body --
    body = communication_data.get('body', '')
    if body:
        paragraphs = body.split('\n\n')
        for para_text in paragraphs:
            para_text = para_text.strip()
            if para_text:
                p = doc.add_paragraph()
                run = p.add_run(para_text)
                run.font.size = Pt(11)

    # -- Call to action --
    cta = communication_data.get('call_to_action', '')
    if cta:
        p = doc.add_paragraph()
        run = p.add_run(cta)
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*accent)

    # -- Closing --
    closing = communication_data.get('closing', '')
    if closing:
        for line in closing.split('\n'):
            line = line.strip()
            if line:
                p = doc.add_paragraph()
                run = p.add_run(line)
                run.font.size = Pt(11)

    # -- Footer --
    footer_para = doc.add_paragraph()
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = footer_para.add_run(brand.get('footer_text', ''))
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(*secondary)

    doc.save(output_path)
    return output_path


def generate_powerpoint(communication_data, brand, output_path):
    """Generate a PowerPoint presentation with brand styling."""
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    primary_hex = '#001D38'
    accent_hex = '#C8A064'
    light_hex = '#ffffff'
    secondary_hex = '#0a2744'
    teal_hex = '#00B398'

    brand_name = brand.get('brand_name', 'Change Management Hub')
    subject = communication_data.get('subject', 'Communication')

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_pptrgb(primary_hex)

    # Brand name top
    txBox = slide.shapes.add_textbox(PptInches(0.8), PptInches(0.5), PptInches(11), PptInches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = brand_name.upper()
    p.font.size = PptPt(12)
    p.font.color.rgb = _hex_to_pptrgb(accent_hex)
    p.font.bold = True

    # Title
    txBox2 = slide.shapes.add_textbox(PptInches(0.8), PptInches(2.5), PptInches(11), PptInches(2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subject
    p2.font.size = PptPt(36)
    p2.font.color.rgb = _hex_to_pptrgb(light_hex)
    p2.font.bold = True

    # Subtitle
    txBox3 = slide.shapes.add_textbox(PptInches(0.8), PptInches(4.5), PptInches(11), PptInches(1))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = f"Change Management Communication — {communication_data.get('tone', 'professional').title()}"
    p3.font.size = PptPt(16)
    p3.font.color.rgb = _hex_to_pptrgb(accent_hex)

    # Accent line
    txBox4 = slide.shapes.add_textbox(PptInches(0.8), PptInches(4.2), PptInches(3), PptInches(0.08))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = '_' * 40
    p4.font.size = PptPt(8)
    p4.font.color.rgb = _hex_to_pptrgb(accent_hex)

    # --- Slide 2: Content ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = slide2.background
    fill2 = bg2.fill
    fill2.solid()
    fill2.fore_color.rgb = _hex_to_pptrgb(secondary_hex)

    # Greeting
    greeting = communication_data.get('greeting', '')
    if greeting:
        tbox = slide2.shapes.add_textbox(PptInches(0.8), PptInches(0.5), PptInches(11), PptInches(0.7))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = greeting
        p.font.size = PptPt(20)
        p.font.color.rgb = _hex_to_pptrgb(light_hex)
        p.font.bold = True

    # Body
    body = communication_data.get('body', '')
    if body:
        body_parts = body.split('\n\n')
        y_pos = 1.5
        for part in body_parts:
            part = part.strip()
            if part:
                tbox = slide2.shapes.add_textbox(
                    PptInches(0.8), PptInches(y_pos), PptInches(11), PptInches(1.2)
                )
                tf = tbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = part
                p.font.size = PptPt(16)
                p.font.color.rgb = _hex_to_pptrgb(light_hex)
                p.space_after = PptPt(12)
                y_pos += 1.3

    # Call to action accent box
    cta = communication_data.get('call_to_action', '')
    if cta:
        tbox = slide2.shapes.add_textbox(PptInches(0.8), PptInches(y_pos + 0.5), PptInches(11), PptInches(0.8))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cta
        p.font.size = PptPt(16)
        p.font.color.rgb = _hex_to_pptrgb(accent_hex)
        p.font.bold = True

    # Date footer
    txBoxF = slide2.shapes.add_textbox(PptInches(0.8), PptInches(6.8), PptInches(5), PptInches(0.4))
    tfF = txBoxF.text_frame
    pF = tfF.paragraphs[0]
    pF.text = brand.get('footer_text', '')
    pF.font.size = PptPt(10)
    pF.font.color.rgb = PptRGBColor(136, 146, 160)

    prs.save(output_path)
    return output_path
